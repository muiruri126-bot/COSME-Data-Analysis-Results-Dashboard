"""
Generate COSME Outputs Progress PowerPoint Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MED_GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)

# Outcome colors
C_MANGROVE  = RGBColor(0x2E, 0x7D, 0x32)
C_SEAWEED   = RGBColor(0x00, 0x89, 0x7B)
C_FOREST    = RGBColor(0x33, 0x69, 0x1E)
C_KNOWLEDGE = RGBColor(0x15, 0x65, 0xC0)
C_RESILIENCE= RGBColor(0xEF, 0x6C, 0x00)
C_COMMUNITY = RGBColor(0x6A, 0x1B, 0x9A)
C_SCHOOLS   = RGBColor(0xF5, 0x7F, 0x17)
C_GOVERNANCE= RGBColor(0xAD, 0x14, 0x57)
C_WRO_YLO   = RGBColor(0x45, 0x27, 0xA0)

C_COMPLETED = RGBColor(0x2E, 0x7D, 0x32)
C_ONGOING   = RGBColor(0xEF, 0x6C, 0x00)
C_NOT_START = RGBColor(0xC6, 0x28, 0x28)
C_TARGET    = RGBColor(0x90, 0xA4, 0xAE)

BRAND_PURPLE = RGBColor(0x4A, 0x14, 0x8C)
BRAND_TEAL   = RGBColor(0x00, 0x89, 0x7B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height


# ═══ Helper functions ═══════════════════════════════════════
def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_kpi_card(slide, left, top, width, height, value, label, color):
    card = add_shape_rect(slide, left, top, width, height, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
    card.shadow.inherit = False
    # Accent bar
    add_shape_rect(slide, left, top, width, Pt(5), color)
    # Value
    add_textbox(slide, left, top + Inches(0.25), width, Inches(0.6), str(value), font_size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, left + Inches(0.1), top + Inches(0.85), width - Inches(0.2), Inches(0.5), label, font_size=11, color=MED_GREY, align=PP_ALIGN.CENTER)
    return card

def make_table(slide, left, top, width, col_widths, headers, rows, header_color=BRAND_PURPLE):
    """Create a styled table on the slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(0.35 * num_rows))
    tbl = table_shape.table

    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF3, 0xE5, 0xF5)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = 'Calibri'
                p.font.color.rgb = DARK_GREY
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape

def add_bar_chart(slide, left, top, width, height, categories, series_data, chart_title='', colors=None):
    """Add a grouped bar chart. series_data = [(name, [values]), ...]"""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)

    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

    # Apply colors
    if colors:
        for idx, color in enumerate(colors):
            if idx < len(chart.series):
                chart.series[idx].format.fill.solid()
                chart.series[idx].format.fill.fore_color.rgb = color

    # Data labels
    for s in chart.series:
        s.has_data_labels = True
        s.data_labels.font.size = Pt(8)
        s.data_labels.number_format = '#,##0'
        s.data_labels.show_value = True

    chart.value_axis.has_title = False
    chart.category_axis.tick_labels.font.size = Pt(8)

    return chart_frame

def add_horizontal_bar_chart(slide, left, top, width, height, categories, series_data, chart_title='', colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)

    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

    if colors:
        for idx, color in enumerate(colors):
            if idx < len(chart.series):
                chart.series[idx].format.fill.solid()
                chart.series[idx].format.fill.fore_color.rgb = color

    for s in chart.series:
        s.has_data_labels = True
        s.data_labels.font.size = Pt(8)
        s.data_labels.number_format = '#,##0'
        s.data_labels.show_value = True

    chart.category_axis.tick_labels.font.size = Pt(8)

    return chart_frame

def add_doughnut_chart(slide, left, top, width, height, categories, values, colors_list, chart_title=''):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Status', values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)

    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

    plot = chart.plots[0]
    series = plot.series[0]
    for i, c in enumerate(colors_list):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c

    series.has_data_labels = True
    series.data_labels.font.size = Pt(10)
    series.data_labels.font.bold = True
    series.data_labels.number_format = '0"%"'
    series.data_labels.show_percentage = False
    series.data_labels.show_value = True

    return chart_frame


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, RGBColor(0x4A, 0x14, 0x8C))

# Title stripe
add_shape_rect(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(3.2), RGBColor(0xFF, 0xFF, 0xFF))

add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1),
            'COSME PROJECT', font_size=18, bold=True, color=BRAND_TEAL, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.2),
            'Outputs Progress Dashboard', font_size=40, bold=True, color=BRAND_PURPLE, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.6),
            'Complete Project Outputs & Deliverables Tracker across all Immediate Outcomes',
            font_size=14, color=MED_GREY, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
            'March 2026', font_size=12, color=MED_GREY, align=PP_ALIGN.CENTER)

# Accent line
add_shape_rect(slide, Inches(5), Inches(3.45), Inches(3.3), Pt(3), BRAND_TEAL)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — Project Summary & KPIs
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
            'Project Overview — Key Performance Indicators', font_size=26, bold=True, color=BRAND_PURPLE)
add_shape_rect(slide, Inches(0.5), Inches(0.85), Inches(4), Pt(3), BRAND_TEAL)

# KPI cards
card_w = Inches(2.8)
card_h = Inches(1.3)
card_y = Inches(1.4)
gap = Inches(0.35)
start_x = Inches(0.6)

add_kpi_card(slide, start_x, card_y, card_w, card_h, '32', 'Total Outputs Tracked', BRAND_PURPLE)
add_kpi_card(slide, start_x + card_w + gap, card_y, card_w, card_h, '18 (56%)', 'Completed / On Track', C_COMPLETED)
add_kpi_card(slide, start_x + 2*(card_w + gap), card_y, card_w, card_h, '10 (31%)', 'Ongoing', C_ONGOING)
add_kpi_card(slide, start_x + 3*(card_w + gap), card_y, card_w, card_h, '4 (13%)', 'Not Started / Behind', C_NOT_START)

# Doughnut chart
add_doughnut_chart(slide, Inches(0.5), Inches(3.2), Inches(5.5), Inches(4),
                   ['Completed (56%)', 'Ongoing (31%)', 'Not Started (13%)'],
                   [18, 10, 4],
                   [C_COMPLETED, C_ONGOING, C_NOT_START],
                   'Output Status Distribution')

# Horizontal bar — Progress by Outcome
outcomes_labels = [
    '1330 WRO/YLO', '1320 Governance', '1310 Schools',
    '1230 Community', '1220 Resilience', '1210 Knowledge',
    '1130 Forest', '1120 Seaweed', '1110 Mangrove'
]
outcomes_vals = [44, 33, 122, 76, 76, 79, 126, 80, 110]
outcome_colors = [C_WRO_YLO, C_GOVERNANCE, C_SCHOOLS, C_COMMUNITY, C_RESILIENCE, C_KNOWLEDGE, C_FOREST, C_SEAWEED, C_MANGROVE]

chart_data = CategoryChartData()
chart_data.categories = outcomes_labels
chart_data.add_series('Avg Progress %', outcomes_vals)

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.5), Inches(3.2), Inches(6.5), Inches(4), chart_data
)
chart = chart_frame.chart
chart.has_legend = False
chart.has_title = True
chart.chart_title.text_frame.paragraphs[0].text = 'Progress by Immediate Outcome (%)'
chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
chart.chart_title.text_frame.paragraphs[0].font.bold = True

series = chart.series[0]
for i, c in enumerate(outcome_colors):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = c

series.has_data_labels = True
series.data_labels.font.size = Pt(9)
series.data_labels.font.bold = True
series.data_labels.number_format = '0"%"'
series.data_labels.show_value = True

chart.category_axis.tick_labels.font.size = Pt(9)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Outcome 1110: Mangrove Restoration
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), C_MANGROVE)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Outcome 1110 — Mangrove Restoration', font_size=26, bold=True, color=WHITE)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
            'Increased capacity of communities, especially women, to undertake gender responsive and equitable mangrove restoration',
            font_size=12, color=MED_GREY)

headers = ['Output', 'Indicator', 'Annual Target', 'Actual (Period)', 'Cumulative', 'End Target', 'Progress', 'Status']
rows = [
    ['1111', 'Biodiversity Assessments conducted', 'N/A', '0', '1', '1', '100%', 'Completed'],
    ['1112a', 'Mangrove Groups established', 'N/A', '3', '42', '40', '105%', 'Completed'],
    ['1112b', 'Members trained (F:840, M:360)', '1,200', '864', '843', '1,200', '70%', 'Ongoing'],
    ['1113', 'Groups equipped with restoration tools', '5', '5', '10', '6', '167%', 'Completed'],
]
col_widths = [Inches(0.7), Inches(2.8), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1), Inches(0.9), Inches(0.9)]
make_table(slide, Inches(0.3), Inches(1.8), Inches(9.7), col_widths, headers, rows, C_MANGROVE)

add_bar_chart(slide, Inches(0.3), Inches(4.0), Inches(6), Inches(3.2),
              ['Biodiversity\nAssessments', 'Groups\nEstablished', 'Members\nTrained', 'Groups\nEquipped'],
              [('End Target', [1, 40, 1200, 6]), ('Cumulative', [1, 42, 843, 10])],
              'Target vs Cumulative', [C_TARGET, C_MANGROVE])

# Status highlights
add_textbox(slide, Inches(6.8), Inches(4.2), Inches(6), Inches(0.4),
            'Key Highlights:', font_size=14, bold=True, color=C_MANGROVE)
add_textbox(slide, Inches(6.8), Inches(4.7), Inches(6), Inches(2.5),
            '• Biodiversity Assessments: 100% completed\n• Mangrove Groups: 42 established vs 40 target (+105%)\n• Groups Equipped: 10 vs 6 target (+167%)\n• Members Training: 843/1,200 (70%) — ongoing',
            font_size=11, color=DARK_GREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — Outcome 1120: Seaweed Production
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), C_SEAWEED)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Outcome 1120 — Seaweed Production', font_size=26, bold=True, color=WHITE)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
            'Increased capacity of women-led seaweed groups to undertake regenerative and sustainable seaweed production, value addition and commercialization',
            font_size=12, color=MED_GREY)

rows = [
    ['1121a', 'Assessments carried out', '2', '0', '2', '2', '100%', 'Completed'],
    ['1121b', 'Seaweed Groups established', '20', '19', '19', '20', '95%', 'Ongoing'],
    ['1122', 'Seaweed Group members', '600', '462', '462', '600', '77%', 'Ongoing'],
    ['1123', 'Seaweed Groups equipped', '20', '19', '19', '20', '95%', 'Ongoing'],
    ['1124', 'Women in value addition (innovation)', '60', '0', '20', '60', '33%', 'Ongoing'],
]
make_table(slide, Inches(0.3), Inches(1.8), Inches(9.7), col_widths, headers, rows, C_SEAWEED)

add_bar_chart(slide, Inches(0.3), Inches(4.3), Inches(6), Inches(3),
              ['Assessments', 'Groups\nEstablished', 'Group\nMembers', 'Groups\nEquipped', 'Women in\nValue Add.'],
              [('End Target', [2, 20, 600, 20, 60]), ('Cumulative', [2, 19, 462, 19, 20])],
              'Target vs Cumulative', [C_TARGET, C_SEAWEED])

add_textbox(slide, Inches(6.8), Inches(4.5), Inches(6), Inches(0.4),
            'Key Highlights:', font_size=14, bold=True, color=C_SEAWEED)
add_textbox(slide, Inches(6.8), Inches(5.0), Inches(6), Inches(2.2),
            '• Assessments: 100% completed\n• Groups established: 19/20 (95%)\n• Group members: 462/600 (77%) — ongoing\n• Groups equipped: 19/20 (95%)\n• Value addition pilots: 20/60 (33%) — needs attention',
            font_size=11, color=DARK_GREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — Outcome 1130: Forest Management
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), C_FOREST)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Outcome 1130 — Forest Management & Conservation', font_size=26, bold=True, color=WHITE)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
            'Increased capacity of communities, especially women, to undertake gender-responsive, locally-led forest management and conservation',
            font_size=12, color=MED_GREY)

rows = [
    ['1131a', 'Authorities sensitized (Kenya + Tanzania)', 'N/A', '—', 'K: F:574; T: T:880', '—', '—', 'Ongoing'],
    ['1131b', 'People trained in forest mgmt', '1,350', 'K:1,030; T:675', 'K: F:1,215; T:1,800', '1,350', '—', 'Ongoing'],
    ['1132', "Women's groups supported", '102', '63', '185', '98', '189%', 'Completed'],
]
make_table(slide, Inches(0.3), Inches(1.8), Inches(9.7), col_widths, headers, rows, C_FOREST)

add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(0.4),
            'Key Highlights:', font_size=14, bold=True, color=C_FOREST)
add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(2.5),
            "• Women's groups supported: 185 vs 98 target (189%) — exceeded target\n"
            '• Training across Kenya and Tanzania ongoing with strong progress\n'
            '• Regional and local authorities sensitized in both countries\n'
            '• Kenya: F:1,215 trained; Tanzania: 1,800 total trained',
            font_size=12, color=DARK_GREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — Outcome 1210: Knowledge & Skills
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), C_KNOWLEDGE)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Outcome 1210 — Knowledge & Skills on Gender Responsive NbS', font_size=26, bold=True, color=WHITE)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
            'Increased knowledge and skills of women on gender responsive NbS, economic rights, life skills, and GE&I, including UPCW',
            font_size=12, color=MED_GREY)

rows = [
    ['1211', 'Women trained on NbS, life skills, GE&I', '3,000', '2,370', '2,370', '3,000', '79%', 'Ongoing'],
]
make_table(slide, Inches(0.3), Inches(1.8), Inches(9.7), col_widths, headers, rows, C_KNOWLEDGE)

# Single big metric
add_kpi_card(slide, Inches(2), Inches(3.5), Inches(3), Inches(1.5), '2,370', 'Women Trained to Date', C_KNOWLEDGE)
add_kpi_card(slide, Inches(5.5), Inches(3.5), Inches(3), Inches(1.5), '3,000', 'End Target', C_TARGET)
add_kpi_card(slide, Inches(9), Inches(3.5), Inches(3), Inches(1.5), '79%', 'Progress', C_ONGOING)

add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5),
            '• 2,370 out of 3,000 women trained (79% of target)\n'
            '• Training covers gender responsive NbS, economic rights, life skills & GE&I\n'
            '• On track to meet target with continued effort',
            font_size=12, color=DARK_GREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — Outcome 1220: Resilience Building
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), C_RESILIENCE)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Outcome 1220 — Resilience Building Assets', font_size=26, bold=True, color=WHITE)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
            'Increased access to resilience building assets and opportunities for women',
            font_size=12, color=MED_GREY)

rows = [
    ['1221a', "Women's Savings Groups established", '130', '—', '140', '130', '108%', 'Completed'],
    ['1221b', 'Savings Groups members trained', '3,000', '2,333', '2,333', '3,000', '78%', 'Ongoing'],
    ['1222a', 'Women trained on regen agriculture', 'F:525, M:225', '—', 'F:525', '750', '70%', 'Ongoing'],
    ['1222b', 'Demonstration plots supported', '14', '0', '0', '14', '0%', 'Ongoing'],
    ['1223', 'Women received solar/tech solutions', '3,000', '—', '3,668', '3,000', '122%', 'Completed'],
]
make_table(slide, Inches(0.3), Inches(1.8), Inches(9.7), col_widths, headers, rows, C_RESILIENCE)

add_bar_chart(slide, Inches(0.3), Inches(4.3), Inches(6), Inches(3),
              ['Savings\nGroups', 'SG\nMembers', 'Regen Ag\nTrained', 'Demo\nPlots', 'Solar/Tech\nSolutions'],
              [('End Target', [130, 3000, 750, 14, 3000]), ('Cumulative', [140, 2333, 525, 0, 3668])],
              'Target vs Cumulative', [C_TARGET, C_RESILIENCE])

add_textbox(slide, Inches(6.8), Inches(4.5), Inches(6), Inches(0.4),
            'Key Highlights:', font_size=14, bold=True, color=C_RESILIENCE)
add_textbox(slide, Inches(6.8), Inches(5.0), Inches(6), Inches(2.2),
            '• Savings Groups: 140/130 (108%) — exceeded target\n'
            '• Solar/Tech solutions: 3,668/3,000 (122%) — exceeded\n'
            '• SG members trained: 2,333/3,000 (78%)\n'
            '• Regen agriculture training: 525/750 (70%)\n'
            '• Demonstration plots: 0/14 — needs immediate attention',
            font_size=11, color=DARK_GREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — Outcome 1230: Community Capacity & Leadership
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), C_COMMUNITY)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Outcome 1230 — Community Capacity & Leadership', font_size=26, bold=True, color=WHITE)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
            'Increased capacity of community members and leaders, particularly men, to promote and support gender responsive NbS',
            font_size=12, color=MED_GREY)

rows = [
    ['1231a', 'SBCC strategy developed', 'N/A', '1', '1', '1', '100%', 'Completed'],
    ['1231b', 'SBCC initiatives undertaken', '93', '75', '79', '90', '88%', 'Ongoing'],
    ['1232a', 'Change Agents trained (F:13, M:11)', 'N/A', 'F:13, M:11', '24', '24', '100%', 'Completed'],
    ['1232b', 'Community sensitization action plans', '24', '24', '24', '24', '100%', 'Completed'],
    ['1233', 'Men and boys trained', '850', '587', '587', '850', '69%', 'Ongoing'],
    ['1234', 'Transformative HH action plans', '850', '0', '0', '850', '0%', 'Ongoing'],
]
make_table(slide, Inches(0.3), Inches(1.75), Inches(9.7), col_widths, headers, rows, C_COMMUNITY)

add_bar_chart(slide, Inches(0.3), Inches(4.5), Inches(6), Inches(2.8),
              ['SBCC\nStrategy', 'SBCC\nInitiatives', 'Change\nAgents', 'Action\nPlans', 'Men/Boys\nTrained', 'HH Action\nPlans'],
              [('End Target', [1, 90, 24, 24, 850, 850]), ('Cumulative', [1, 79, 24, 24, 587, 0])],
              'Target vs Cumulative', [C_TARGET, C_COMMUNITY])

add_textbox(slide, Inches(6.8), Inches(4.7), Inches(6), Inches(2.2),
            '• SBCC strategy, Change Agents, Action Plans: all 100%\n'
            '• SBCC initiatives: 79/90 (88%) — nearly complete\n'
            '• Men & boys trained: 587/850 (69%) — ongoing\n'
            '• HH action plans: 0/850 — not yet started',
            font_size=11, color=DARK_GREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — Outcome 1310: School Children Awareness
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), C_SCHOOLS)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Outcome 1310 — School Children Awareness', font_size=26, bold=True, color=WHITE)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
            'Increased awareness and knowledge of primary school children, particularly girls, on climate change, NbS, and conservation',
            font_size=12, color=MED_GREY)

rows = [
    ['1311a', '4K and Roots & Shoots clubs established', '110', 'K:66, T:30', 'K:66, T:45', '110', '101%', 'Completed'],
    ['1311b', '4K & R&S club members trained', '6,150', '11,815', '12,175', 'F:3,661 M:2,489', '198%', 'Completed'],
    ['1312', 'Clubs implementing conservation', '110', 'K:12, T:45', '98', '110', '89%', 'Ongoing'],
    ['1313', 'Schools receiving clean water', '65', '65', '65', '65', '100%', 'Completed'],
]
make_table(slide, Inches(0.3), Inches(1.8), Inches(9.7), col_widths, headers, rows, C_SCHOOLS)

add_bar_chart(slide, Inches(0.3), Inches(4.0), Inches(6), Inches(3.2),
              ['Clubs\nEstablished', 'Members\nTrained', 'Clubs\nConserving', 'Clean Water\nSchools'],
              [('End Target', [110, 6150, 110, 65]), ('Cumulative', [111, 12175, 98, 65])],
              'Target vs Cumulative', [C_TARGET, C_SCHOOLS])

add_textbox(slide, Inches(6.8), Inches(4.2), Inches(6), Inches(0.4),
            'Key Highlights:', font_size=14, bold=True, color=C_SCHOOLS)
add_textbox(slide, Inches(6.8), Inches(4.7), Inches(6), Inches(2.5),
            '• Club members trained: 12,175 vs 6,150 target (198%)\n'
            '• 4K/R&S clubs: 111 established vs 110 target (101%)\n'
            '• Clean water solutions: 65/65 (100%)\n'
            '• Conservation clubs: 98/110 (89%) — nearly there',
            font_size=11, color=DARK_GREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — Outcome 1320: Community Governance
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), C_GOVERNANCE)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Outcome 1320 — Community Governance', font_size=26, bold=True, color=WHITE)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
            'Strengthened gender responsive community governance structures to reduce risk and enhance preparedness to climate change',
            font_size=12, color=MED_GREY)

rows = [
    ['1321', 'Gender responsive adaptation & DRR plans', '5', '0', '0', '65', '0%', 'Ongoing'],
    ['1322', 'Communities supported', 'N/A', '0', '0', '30', '0%', 'Not Started'],
    ['1323', 'Learning and sharing events held', 'N/A', '0', '1', '1', '100%', 'Completed'],
]
make_table(slide, Inches(0.3), Inches(1.8), Inches(9.7), col_widths, headers, rows, C_GOVERNANCE)

# KPI cards for this small outcome
add_kpi_card(slide, Inches(1.5), Inches(3.8), Inches(3), Inches(1.5), '0/65', 'Adaptation & DRR Plans', C_NOT_START)
add_kpi_card(slide, Inches(5.2), Inches(3.8), Inches(3), Inches(1.5), '0/30', 'Communities Supported', C_NOT_START)
add_kpi_card(slide, Inches(8.9), Inches(3.8), Inches(3), Inches(1.5), '1/1', 'Learning Events', C_COMPLETED)

add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(1.5),
            '⚠ This outcome has the lowest progress (33%). Adaptation/DRR plans and community support activities need '
            'immediate prioritization. The single learning event has been completed.',
            font_size=13, color=C_NOT_START, bold=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 11 — Outcome 1330: WRO & YLO Advocacy
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), C_WRO_YLO)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Outcome 1330 — WRO & YLO Advocacy', font_size=26, bold=True, color=WHITE)

add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
            'Increased ability of WRO and YLOs to undertake evidenced-based advocacy for gender responsive and inclusive climate adaptation',
            font_size=12, color=MED_GREY)

rows = [
    ['1331', 'Evidence on drivers of gender responsive adaptation', 'N/A', '0', '0', '7', '0%', 'Ongoing'],
    ['1332', 'Dissemination events', '16', '1', '0', '18', '0%', 'Ongoing'],
    ['1333', 'WRO/YLO members trained (F:150, M:50)', '200', 'F:150', '—', '200', '75%', 'Ongoing'],
    ['1334', 'WRO/YLO evidence-based advocacy plans', '4', '0', '4', '4', '100%', 'Not Started'],
]
make_table(slide, Inches(0.3), Inches(1.8), Inches(9.7), col_widths, headers, rows, C_WRO_YLO)

add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(0.4),
            'Key Highlights:', font_size=14, bold=True, color=C_WRO_YLO)
add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(2.5),
            '• WRO/YLO training: 150 females trained, 50 males planned (75% progress)\n'
            '• Evidence-based advocacy plans: 4/4 completed (100%)\n'
            '• Evidence gathering on drivers: 0/7 — needs acceleration\n'
            '• Dissemination events: 0/18 — significant gap to close',
            font_size=12, color=DARK_GREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 12 — Gender Breakdown in Training
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), BRAND_PURPLE)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Cross-Cutting Analysis — Gender in Training Outputs', font_size=26, bold=True, color=WHITE)

# Female vs Male training bar chart
add_bar_chart(slide, Inches(0.3), Inches(1.4), Inches(8), Inches(5.5),
              ['Mangrove\nTraining', 'Forest Mgmt\n(Kenya)', 'Women NbS\nTraining', 'Regen Ag\nTraining', 'Change\nAgents', 'WRO/YLO\nTraining'],
              [('Female', [849, 1215, 2370, 525, 13, 150]), ('Male', [194, 585, 0, 225, 11, 50])],
              'Female vs Male — Training Outputs', [RGBColor(0xAB, 0x47, 0xBC), BRAND_TEAL])

# 4K Club members by country
add_bar_chart(slide, Inches(8.5), Inches(1.4), Inches(4.5), Inches(5.5),
              ['Kenya', 'Tanzania'],
              [('Female', [5965, 6810]), ('Male', [3217, 3464])],
              '4K/R&S Club Members by Gender', [RGBColor(0xAB, 0x47, 0xBC), BRAND_TEAL])


# ═══════════════════════════════════════════════════════════
# SLIDE 13 — Summary & Recommendations
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), SW, Inches(1.05), BRAND_PURPLE)
add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
            'Summary & Recommendations', font_size=26, bold=True, color=WHITE)

# Achievements column
add_shape_rect(slide, Inches(0.4), Inches(1.3), Inches(6), Inches(5.8), WHITE, RGBColor(0xE0, 0xE0, 0xE0))
add_shape_rect(slide, Inches(0.4), Inches(1.3), Inches(6), Pt(5), C_COMPLETED)
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.5),
            '✓  Key Achievements', font_size=18, bold=True, color=C_COMPLETED)
add_textbox(slide, Inches(0.7), Inches(2.1), Inches(5.5), Inches(4.8),
            '• 18 of 32 outputs completed or on track (56%)\n\n'
            '• Mangrove groups exceeded target: 42/40 (105%)\n\n'
            '• Solar/tech solutions exceeded: 3,668/3,000 (122%)\n\n'
            '• School club members: 12,175 vs 6,150 target (198%)\n\n'
            "• Women's groups supported: 185 vs 98 target (189%)\n\n"
            '• Savings groups exceeded target: 140/130 (108%)\n\n'
            '• All Change Agents trained, all action plans completed',
            font_size=12, color=DARK_GREY)

# Attention column
add_shape_rect(slide, Inches(6.9), Inches(1.3), Inches(6), Inches(5.8), WHITE, RGBColor(0xE0, 0xE0, 0xE0))
add_shape_rect(slide, Inches(6.9), Inches(1.3), Inches(6), Pt(5), C_NOT_START)
add_textbox(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.5),
            '⚠  Areas Needing Attention', font_size=18, bold=True, color=C_NOT_START)
add_textbox(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(4.8),
            '• Governance (1320) at 33% — DRR plans 0/65\n\n'
            '• WRO/YLO advocacy (1330) at 44% — events 0/18\n\n'
            '• Demonstration plots: 0/14 — not started\n\n'
            '• HH action plans: 0/850 — not yet started\n\n'
            '• Seaweed value addition: 20/60 (33%)\n\n'
            '• Men & boys training: 587/850 (69%)\n\n'
            '• Dissemination events and evidence generation behind',
            font_size=12, color=DARK_GREY)


# ═══ Save ═══════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'COSME_Outputs_Progress_Presentation.pptx')
prs.save(output_path)
print(f'✅ Presentation saved to: {output_path}')
print(f'   Slides: {len(prs.slides)}')
